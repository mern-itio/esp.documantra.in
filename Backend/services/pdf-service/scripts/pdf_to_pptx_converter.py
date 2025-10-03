#!/usr/bin/env python3
import json
import os
import sys
import math
import tempfile

try:
    import fitz  # PyMuPDF
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from PIL import Image
except Exception as e:
    print(json.dumps({"success": False, "error": f"Missing dependency: {e}"}))
    sys.exit(1)


def px_to_inches(px, dpi):
    return px / float(dpi)


def add_image(slide, image_path, left_in, top_in, width_in, height_in):
    slide.shapes.add_picture(image_path, Inches(left_in), Inches(top_in), width=Inches(width_in), height=Inches(height_in))


def add_textbox(slide, text, left_in, top_in, width_in, height_in, font_size_pt=12):
    tx_box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    tf = tx_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(font_size_pt)
    font.color.rgb = RGBColor(0, 0, 0)


def main():
    if len(sys.argv) < 3:
        print(json.dumps({"success": False, "error": "Usage: pdf_to_pptx_converter.py <input.pdf> <output.pptx>"}))
        return 1

    input_pdf = os.path.abspath(sys.argv[1])
    output_pptx = os.path.abspath(sys.argv[2])

    if not os.path.exists(input_pdf):
        print(json.dumps({"success": False, "error": f"Input PDF not found: {input_pdf}"}))
        return 1

    prs = Presentation()
    doc = fitz.open(input_pdf)

    # Render images to temp dir for placement when needed
    tmp_dir = tempfile.mkdtemp(prefix="pdf2pptx_")
    try:
        for page_index in range(len(doc)):
            page = doc[page_index]
            # Compute slide size to match page aspect ratio and avoid distortion
            pdf_width_pt, pdf_height_pt = page.rect.width, page.rect.height
            pdf_width_in = pdf_width_pt / 72.0
            pdf_height_in = pdf_height_pt / 72.0

            # Target slide height 7.5in, compute width by aspect; clamp width to common range
            target_height_in = 7.5
            computed_width_in = max(8.0, min(13.33, target_height_in * (pdf_width_in / pdf_height_in)))
            prs.slide_height = Inches(target_height_in)
            prs.slide_width = Inches(computed_width_in)
            slide_width_in = computed_width_in
            slide_height_in = target_height_in

            scale_x = slide_width_in / pdf_width_in
            scale_y = slide_height_in / pdf_height_in

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

            # Extract text with rich structure to better preserve positions
            text_dict = page.get_text("rawdict")

            # Place images at their actual rectangles
            images = page.get_images(full=True)
            for img_index, img in enumerate(images):
                xref = img[0]
                try:
                    rects = page.get_image_rects(xref)
                except Exception:
                    rects = []
                base_image = doc.extract_image(xref)
                img_bytes = base_image.get("image")
                ext = base_image.get("ext", "png")
                img_path = os.path.join(tmp_dir, f"p{page_index}_img{img_index}.{ext}")
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                if rects:
                    for ridx, r in enumerate(rects):
                        left_in = (r.x0 / 72.0) * scale_x
                        top_in = (r.y0 / 72.0) * scale_y
                        width_in = ((r.x1 - r.x0) / 72.0) * scale_x
                        height_in = ((r.y1 - r.y0) / 72.0) * scale_y
                        add_image(slide, img_path, left_in, top_in, width_in, height_in)
                else:
                    # Fallback: center fit if rect unknown
                    try:
                        with Image.open(img_path) as im:
                            img_w_in = im.width / 96.0
                            img_h_in = im.height / 96.0
                    except Exception:
                        img_w_in, img_h_in = slide_width_in, slide_height_in
                    ratio = min(slide_width_in / img_w_in, slide_height_in / img_h_in)
                    w_in = img_w_in * ratio
                    h_in = img_h_in * ratio
                    left_in = (slide_width_in - w_in) / 2.0
                    top_in = (slide_height_in - h_in) / 2.0
                    add_image(slide, img_path, left_in, top_in, w_in, h_in)

            placed_any_text = False
            # Add text using spans to reduce overlap and better font sizing
            try:
                for block in text_dict.get("blocks", []):
                    if block.get("type") != 0:
                        continue
                    for line in block.get("lines", []):
                        # Compute bbox for the line
                        bbox = None
                        line_text = ""
                        max_size = 12
                        for span in line.get("spans", []):
                            line_text += span.get("text", "")
                            size = span.get("size")
                            if size:
                                max_size = max(max_size, int(size))
                            if not bbox:
                                bbox = list(span.get("bbox", [0, 0, 0, 0]))
                            else:
                                sb = span.get("bbox", [0, 0, 0, 0])
                                bbox[0] = min(bbox[0], sb[0])
                                bbox[1] = min(bbox[1], sb[1])
                                bbox[2] = max(bbox[2], sb[2])
                                bbox[3] = max(bbox[3], sb[3])
                        if not line_text.strip() or not bbox:
                            continue
                        x0, y0, x1, y1 = bbox
                        left_in = (x0 / 72.0) * scale_x
                        top_in = (y0 / 72.0) * scale_y
                        width_in = max(0.5, ((x1 - x0) / 72.0) * scale_x)
                        height_in = max(0.3, ((y1 - y0) / 72.0) * scale_y)
                        font_size = max(8, int(max_size))
                        add_textbox(slide, line_text.strip(), left_in, top_in, width_in, height_in, font_size_pt=font_size)
                        placed_any_text = True
            except Exception:
                # Fallback: simple blocks
                text_blocks = page.get_text("blocks")
                for block in text_blocks:
                    if len(block) < 5:
                        continue
                    x0, y0, x1, y1, text = block[:5]
                    if not text or not text.strip():
                        continue
                    left_in = (x0 / 72.0) * scale_x
                    top_in = (y0 / 72.0) * scale_y
                    width_in = max(0.5, ((x1 - x0) / 72.0) * scale_x)
                    height_in = max(0.3, ((y1 - y0) / 72.0) * scale_y)
                    font_size = max(8, int((y1 - y0) / 72.0 * 10))
                    add_textbox(slide, text.strip(), left_in, top_in, width_in, height_in, font_size_pt=font_size)
                    placed_any_text = True

            # If no text placed yet, fallback to word-level extraction
            if not placed_any_text:
                try:
                    words = page.get_text("words")  # x0, y0, x1, y1, word, block_no, line_no, word_no
                    from collections import defaultdict
                    lines = defaultdict(list)
                    for w in words:
                        if len(w) < 8:
                            continue
                        x0, y0, x1, y1, word, block_no, line_no, word_no = w
                        lines[(block_no, line_no)].append((word_no, x0, y0, x1, y1, word))

                    for (_b, _l), wl in lines.items():
                        wl.sort(key=lambda t: t[0])
                        if not wl:
                            continue
                        text = " ".join([t[5] for t in wl]).strip()
                        if not text:
                            continue
                        x0 = min(t[1] for t in wl)
                        y0 = min(t[2] for t in wl)
                        x1 = max(t[3] for t in wl)
                        y1 = max(t[4] for t in wl)
                        left_in = (x0 / 72.0) * scale_x
                        top_in = (y0 / 72.0) * scale_y
                        width_in = max(0.5, ((x1 - x0) / 72.0) * scale_x)
                        height_in = max(0.3, ((y1 - y0) / 72.0) * scale_y)
                        # Estimate font size from line height
                        font_size = max(8, int((y1 - y0) / 72.0 * 72 * 0.35))
                        add_textbox(slide, text, left_in, top_in, width_in, height_in, font_size_pt=font_size)
                        placed_any_text = True
                except Exception:
                    pass

        prs.save(output_pptx)
        size = os.path.getsize(output_pptx) if os.path.exists(output_pptx) else 0
        print(json.dumps({"success": True, "output": output_pptx, "fileSize": size, "pages": len(doc)}))
        return 0
    except Exception as e:
        print(json.dumps({"success": False, "error": str(e)}))
        return 1
    finally:
        try:
            import shutil
            shutil.rmtree(tmp_dir)
        except Exception:
            pass


if __name__ == "__main__":
    sys.exit(main())


